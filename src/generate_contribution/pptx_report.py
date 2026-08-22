"""PowerPoint report assembly: builds one slide group per indicator (descriptive
table, winsorization/Box-Cox summaries, normalized scatter + normality tests),
each with a boxplot/histogram chart and a choropleth map, plus a title slide
and a sector diagram slide.

One behavior worth knowing about: `slides_resultT` assumes `resumo.resumo_basico`'s columns
and `datawinz.resumo` / `databxcx.meta` / `datanorm.idata`'s rows/columns line up 1:1 per
indicator by position. That holds whenever `iMeta` has no "Cluster"-classified rows. A
"Cluster"-classified indicator legitimately produces 3 descriptive views in `resumo_basico`
(Conjunto Completo/Grupo 1/Grupo 2) but only 2 winsorization rows in `datawinz.resumo` (Grupo
1/Grupo 2 -- see `winsorise.py`), so the 1:1 assumption doesn't hold for such datasets. Rather
than silently misindexing, this raises a clear error in that case; per-indicator slide layout
for Cluster-classified indicators is unimplemented.
"""
from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
import shutil
import tempfile
import warnings

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt
from scipy import stats as sp_stats

from .boxcox import BoxCoxResult
from .config import ReportConfig
from .correlation import CorrelIndResult
from .diagrams import gerar_diagrama_setor
from .figures import criar_grafico, grafico_final
from .maps import Map_result, load_shapefiles, mapnorma_result
from .normalise import NormaliseResult
from .resumo import ResumoResult
from .winsorise import WinsoriseResult

_TITLE_AND_CONTENT_LAYOUT = "Title and Content"


def _get_layout(prs: Presentation, name: str = _TITLE_AND_CONTENT_LAYOUT):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[1]  # conventional index for "Title and Content" in the default template


def _set_title(slide, text: str, size: int = 22) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size)


def _add_textbox(slide, text: str, left_in: float, top_in: float, size: int, color: tuple[int, int, int], bold: bool = True) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(8), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _add_table(
    slide,
    df: pd.DataFrame,
    left_in: float,
    top_in: float,
    col_widths_in: list[float],
    row_height_in: float = 0.35,
    header_height_in: float = 0.5,
    header_align: PP_ALIGN | None = None,
    body_align: dict[int, PP_ALIGN] | None = None,
    merge_header: bool = False,
):
    """Adds a table with a bold header row, per-column widths, and optional cell alignment."""
    body_align = body_align or {}
    n_rows, n_cols = df.shape
    width_total = Inches(sum(col_widths_in))
    height_total = Inches(header_height_in + row_height_in * n_rows)
    shape = slide.shapes.add_table(n_rows + 1, n_cols, Inches(left_in), Inches(top_in), width_total, height_total)
    table = shape.table

    for j, w in enumerate(col_widths_in):
        table.columns[j].width = Inches(w)
    table.rows[0].height = Inches(header_height_in)
    for i in range(n_rows):
        table.rows[i + 1].height = Inches(row_height_in)

    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        if header_align is not None:
            cell.text_frame.paragraphs[0].alignment = header_align

    if merge_header and n_cols > 1:
        table.cell(0, 0).merge(table.cell(0, n_cols - 1))

    for i in range(n_rows):
        for j in range(n_cols):
            cell = table.cell(i + 1, j)
            cell.text = "" if pd.isna(df.iat[i, j]) else str(df.iat[i, j])
            if j in body_align:
                cell.text_frame.paragraphs[0].alignment = body_align[j]

    return table


def _add_picture(slide, path: str | Path, left_in: float, top_in: float, width_cm: float, height_cm: float) -> None:
    if path is None or not Path(path).exists():
        return
    slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), width=Cm(width_cm), height=Cm(height_cm))


@dataclass(slots=True)
class PptxContext:
    prs: Presentation
    outfile: Path


def create_pptx(
    template: str | Path,
    meta_dados: pd.DataFrame,
    setor_estrategico: str,
    sigla: str,
    subsetor: str | None,
    output_dir: Path,
    temp_dir: Path,
) -> PptxContext:
    """Copies the template, sets the title slide's subtitle, and adds the sector diagram slide."""
    output_dir.mkdir(parents=True, exist_ok=True)
    suffix = sigla + (subsetor or "")
    ts = datetime.now().strftime("%d-%m-%Y_%Hh%Mm")
    outfile = output_dir / f"REL_{suffix}_{ts}.pptx"
    sub_title = f"{setor_estrategico}-{subsetor}" if subsetor else setor_estrategico

    shutil.copy(template, outfile)
    prs = Presentation(str(outfile))

    slide1 = prs.slides[0]
    caixa = next((s for s in slide1.shapes if s.shape_id == 6), None)
    if caixa is not None:
        left_in = caixa.left / Inches(1)
        top_in = caixa.top / Inches(1)
        width_in = caixa.width / Inches(1)
        height_in = caixa.height / Inches(1)
        caixa.text_frame.clear()
        p = caixa.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = sub_title
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.italic = True
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        caixa.left = Inches(left_in)
        caixa.top = Inches(top_in + 0.8)
        caixa.width = Inches(width_in)
        caixa.height = Inches(height_in)

    diagrama_path = temp_dir / "diagrama.png"
    gerar_diagrama_setor(meta_dados, sub_title, diagrama_path, width=1600, height=1000)

    slide2 = prs.slides.add_slide(_get_layout(prs))
    _set_title(slide2, "Indicadores e Indices")
    _add_picture(slide2, diagrama_path, left_in=0.5, top_in=0.75, width_cm=31.6, height_cm=14)

    return PptxContext(prs=prs, outfile=outfile)


def slides_descricao(
    ctx: PptxContext,
    tabela1: pd.DataFrame,
    tabela2: pd.DataFrame,
    tabela3: pd.DataFrame,
    titulo: str,
    caminho_imagem: str | Path | None = None,
    caminho_map: str | Path | None = None,
) -> None:
    """Adds the descriptive-statistics slide: metadata table, stats table, NA/uniqueness table, chart + map."""
    slide = ctx.prs.slides.add_slide(_get_layout(ctx.prs))
    _set_title(slide, titulo)

    if tabela2["Descritivo"].iloc[0] not in ("Grupo 1", "Grupo 2"):
        _add_table(slide, tabela1, left_in=1, top_in=1.5, col_widths_in=[0.75, 0.75, 1, 5.0, 0.8, 1.4])
        _add_textbox(slide, "Resumo Descritivo", 1, -0.2, size=14, color=(0, 0, 0))

    _add_table(
        slide,
        tabela2,
        left_in=1,
        top_in=2.5,
        col_widths_in=[1.45, 1.45],
        header_height_in=0.5,
        row_height_in=0.35,
        merge_header=True,
        body_align={1: PP_ALIGN.RIGHT},
    )
    _add_table(
        slide,
        tabela3,
        left_in=1,
        top_in=5.9,
        col_widths_in=[1.65] * tabela3.shape[1],
        body_align={j: PP_ALIGN.RIGHT for j in range(tabela3.shape[1])},
    )
    _add_textbox(slide, "Avaliação de NA e Valores Únicos", 1.8, 4.2, size=13, color=(0, 0, 0))

    _add_picture(slide, caminho_imagem, left_in=4.2, top_in=2.6, width_cm=12.5, height_cm=7.40)
    _add_picture(slide, caminho_map, left_in=9.5, top_in=2.06, width_cm=9.04, height_cm=10.85)


def slides_generico(
    ctx: PptxContext,
    tabela1: pd.DataFrame | None,
    titulo: str,
    caminho_map: str | Path | None = None,
) -> None:
    """Adds a generic summary-table + map slide (used for the winsorization/Box-Cox summaries)."""
    slide = ctx.prs.slides.add_slide(_get_layout(ctx.prs))
    _set_title(slide, titulo)

    if tabela1 is not None:
        widths = [1.4 if len(str(c)) <= 10 else 2.5 for c in tabela1.columns]
        _add_table(slide, tabela1, left_in=1, top_in=1.5, col_widths_in=widths)
        _add_textbox(slide, "Resumo", 1, -0.2, size=14, color=(0, 0, 0))

    _add_picture(slide, caminho_map, left_in=4.55, top_in=2.1, width_cm=11.08, height_cm=13.5)


def slides_normal(
    ctx: PptxContext,
    titulo: str,
    dtabxcx: pd.Series,
    caminho_grafico: str | Path | None = None,
    caminho_map: str | Path | None = None,
) -> None:
    """Adds the normalized-data slide: Shapiro-Wilk/Anderson-Darling normality tests, scatter + map."""
    values = pd.to_numeric(pd.Series(dtabxcx), errors="coerce").dropna().to_numpy()

    try:
        with warnings.catch_warnings():
            # deliberately using the legacy critical_values-based result (no
            # `method=`); scipy >=1.17 warns about its future default either way
            warnings.simplefilter("ignore", FutureWarning)
            ad = sp_stats.anderson(values, dist="norm")
        significant = ad.statistic > ad.critical_values[2]  # 5% significance level
        adtest_txt = (
            f"Anderson-Darling normality test\nA = {ad.statistic:.4f}, "
            f"{'rejects' if significant else 'does not reject'} normality at 5%"
        )
    except Exception:
        adtest_txt = "Series de dados incompativel com o AD.test"

    try:
        stat, pvalue = sp_stats.shapiro(values)
        shtest_txt = f"Shapiro-Wilk normality test\nW = {stat:.4f}, p-value = {pvalue:.4g}"
    except Exception:
        shtest_txt = "Series de dados incompativel com o Shapiro.test"

    slide = ctx.prs.slides.add_slide(_get_layout(ctx.prs))
    _set_title(slide, titulo)
    _add_textbox(slide, shtest_txt, 1, -0.1, size=15, color=(0, 0, 0x8B))
    _add_textbox(slide, adtest_txt, 1, 0.5, size=15, color=(0x8B, 0, 0))

    _add_picture(slide, caminho_map, left_in=7.332, top_in=0.28, width_cm=15.08, height_cm=18.5)
    _add_picture(slide, caminho_grafico, left_in=0.176, top_in=3.012, width_cm=17.75, height_cm=8.87)


def slides_resultT(
    template: str | Path,
    setor_estrategico: str,
    sigla: str,
    caminho_shp_mun: str | Path,
    caminho_shp_uf: str | Path,
    resumo: ResumoResult,
    meta_adapta: pd.DataFrame,
    reference: pd.DataFrame,
    databruto: pd.DataFrame,
    datawinz: WinsoriseResult,
    databxcx: BoxCoxResult,
    datanorm: NormaliseResult,
    subsetor: str | None = None,
    ind: int | None = None,
    resu: bool = True,
    winz: bool = True,
    bxcx: bool = True,
    norm: bool = True,
    output_dir: Path = Path("OUTPUT"),
) -> Path:
    """Builds the full report: title/diagram slides, then one slide group per indicator."""
    with tempfile.TemporaryDirectory(prefix="generate_contribution_") as tmp:
        temp_dir = Path(tmp)
        ctx = create_pptx(template, meta_adapta, setor_estrategico, sigla, subsetor, Path(output_dir), temp_dir)

        shp_mun, shp_uf = load_shapefiles(caminho_shp_mun, caminho_shp_uf)

        res_basico = resumo.resumo_basico
        winz_resumo = datawinz.resumo
        bxcx_meta = databxcx.meta
        k = ind if ind is not None else res_basico.shape[1]

        if winz and winz_resumo.shape[0] != res_basico.shape[1]:
            raise ValueError(
                "slides_resultT: datawinz.resumo has a different row count than resumo.resumo_basico has "
                "columns -- the per-indicator positional alignment this loop relies on doesn't hold for this "
                "dataset. This happens when iMeta has any 'Cluster'-classified rows: resumo_basico gets 3 "
                "descriptive views per Cluster indicator (Conjunto Completo/Grupo 1/Grupo 2) but "
                "datawinz.resumo only gets 2 (Grupo 1/Grupo 2). Per-indicator slide layout for "
                "Cluster-classified indicators is not implemented; re-check the input metadata."
            )
        if bxcx and bxcx_meta.shape[0] != res_basico.shape[1]:
            raise ValueError(
                "slides_resultT: databxcx.meta has a different row count than resumo.resumo_basico has columns."
            )

        for i in range(k):
            icode = res_basico.columns[i]
            tab_descricao = meta_adapta.loc[meta_adapta["Code"] == icode, meta_adapta.columns != "Classe"]
            tab_stats = pd.DataFrame(
                {"Resumo Estatístico": res_basico.index, "Descritivo": res_basico.iloc[:, i].to_numpy()}
            )
            tab_nas = resumo.resumo_na.iloc[[i], 1:]

            title_slide = f"{tab_descricao['Nome'].iloc[0]} - {icode}"
            idata = databruto[icode]

            titulo_map = [f"Indicador Bruto: {icode}", f"Indicador Winsorization: {icode}", f"Indicador BoxCox: {icode}", f"Indicador Normalizado: {icode}"]

            gra_descricao_path = temp_dir / "gra_descricao.png"
            criar_grafico(idata, gra_descricao_path, largura=24, altura=12, nvalores=icode, fsize=50)

            map_bruto_path = temp_dir / "map_bruto.png"
            Map_result(icode, idata, reference["GEOCOD"], titulo_map[0], shp_mun, shp_uf, map_bruto_path, largura=20, altura=24, fsize=36)

            if resu:
                slides_descricao(ctx, tab_descricao, tab_stats, tab_nas, title_slide, gra_descricao_path, map_bruto_path)

            if winz:
                iwinsor = winz_resumo.iloc[[i]].copy()
                if iwinsor["Aplicacao"].iloc[0] == "N":
                    iwinsor["Aplicacao"] = "Não foi aplicado winsorization (Sem Mapa)"
                    slides_generico(ctx, iwinsor, title_slide, None)
                else:
                    map_wins_path = temp_dir / "map_Wins.png"
                    Map_result(icode, datawinz.idata[icode], reference["GEOCOD"], titulo_map[1], shp_mun, shp_uf, map_wins_path, largura=20, altura=24, fsize=36)
                    aplicacao_txt = {
                        "A": "Winsorization aplicado nos limites inferior e superior",
                        "S": "Winsorization aplicado no limite superior",
                        "I": "Winsorization aplicado no limite inferior",
                    }.get(iwinsor["Aplicacao"].iloc[0], iwinsor["Aplicacao"].iloc[0])
                    iwinsor["Aplicacao"] = aplicacao_txt
                    slides_generico(ctx, iwinsor, title_slide, map_wins_path)

            if bxcx:
                ibxcx = bxcx_meta.iloc[[i]].copy()
                if ibxcx["BoxCox"].iloc[0] == 0:
                    ibxcx["BoxCox"] = "Não foi aplicado BoxCox (Sem Mapa)"
                    slides_generico(ctx, ibxcx, title_slide, None)
                else:
                    map_boxcox_path = temp_dir / "map_Boxcox.png"
                    Map_result(icode, databxcx.data[icode], reference["GEOCOD"], titulo_map[2], shp_mun, shp_uf, map_boxcox_path, largura=20, altura=24, fsize=36)
                    ibxcx["BoxCox"] = "BoxCox aplicado"
                    slides_generico(ctx, ibxcx, title_slide, map_boxcox_path)

            if norm:
                ibxcx_data = databxcx.data[icode]
                norma_data1 = datanorm.idata[[icode]]
                datacombinado = pd.DataFrame({"Normalizado": datanorm.idata[icode], "N_Normalizado": databruto[icode]}).dropna()

                grafico_final_path = temp_dir / "grafico_final.png"
                grafico_final(
                    df1=datacombinado,
                    df=pd.DataFrame({icode: datanorm.idata[icode].dropna()}),
                    nome_arquivo=grafico_final_path,
                    largura=20,
                    altura=10,
                    dpi=100,
                    nvalores=icode,
                    fsize=32,
                )
                map_norma_path = temp_dir / "map_Norma.png"
                mapnorma_result(icode, norma_data1[icode], reference["GEOCOD"], titulo_map[3], shp_mun, shp_uf, map_norma_path, largura=20, altura=24, fsize=36)

                slides_normal(ctx, title_slide, ibxcx_data, grafico_final_path, map_norma_path)

            print(f"{i + 1}  de  {k}")

        ctx.prs.save(str(ctx.outfile))
        return ctx.outfile
