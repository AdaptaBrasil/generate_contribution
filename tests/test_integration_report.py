"""End-to-end integration test against the real report assets vendored in
`DATASET/`/`TEMPLATE/` (template pptx, real shapefiles). Skipped automatically
when those assets or the system Graphviz `dot` binary aren't available (e.g.
in a CI checkout that doesn't fetch the large binary fixtures, or a dev
machine without Graphviz installed).
"""
from __future__ import annotations

import shutil
from pathlib import Path

import numpy as np
import pandas as pd
import pytest
from pptx import Presentation

from generate_contribution.config import TratamentoConfig
from generate_contribution.pptx_report import slides_resultT
from generate_contribution.treatment import run_tratamento

_PROJECT_ROOT = Path(__file__).resolve().parent.parent
_TEMPLATE = _PROJECT_ROOT / "TEMPLATE" / "ADAPTA_RESUMO.pptx"
_SHP_MUN = _PROJECT_ROOT / "DATASET" / "SHP" / "BR_Municipios_2022_gr.shp"
_SHP_UF = _PROJECT_ROOT / "DATASET" / "SHP" / "BR_UF_2022_gr.shp"

pytestmark = pytest.mark.skipif(
    not _TEMPLATE.exists() or shutil.which("dot") is None,
    reason="requires the DATASET/TEMPLATE report assets and a system Graphviz install",
)


def test_slides_resultT_end_to_end(tmp_path):
    rng = np.random.default_rng(7)
    n = 200
    imeta = pd.DataFrame(
        [
            dict(N=1, Nivel=0, Code="GEOCOD", Nome="Geocodigo", Tipo="Referencia", Pai="REF", Classe="Descricao"),
            dict(N=2, Nivel=0, Code="MUN", Nome="Municipio", Tipo="Referencia", Pai="REF", Classe="Descricao"),
            dict(N=3, Nivel=0, Code="UF", Nome="UF", Tipo="Referencia", Pai="REF", Classe="Descricao"),
            dict(N=4, Nivel=0, Code="CLUSTER", Nome="Cluster", Tipo="Referencia", Pai="REF", Classe="Descricao"),
            dict(N=5, Nivel=7, Code="IND1", Nome="Indicador Um", Tipo="Indicator", Pai="PP", Classe="Numérico"),
            dict(N=6, Nivel=7, Code="IND2", Nome="Indicador Dois", Tipo="Indicator", Pai="IB", Classe="Numérico"),
        ]
    )
    idata = pd.DataFrame(
        {
            "GEOCOD": np.arange(1100000, 1100000 + n),
            "MUN": [f"Municipio {i % 50}" for i in range(n)],
            "UF": ["RO"] * n,
            "CLUSTER": rng.integers(1, 3, size=n),
            "IND1": rng.normal(10, 2, size=n),
            "IND2": rng.uniform(1, 5, size=n),
        }
    )
    input_path = tmp_path / "input.xlsx"
    with pd.ExcelWriter(input_path, engine="openpyxl") as writer:
        imeta.to_excel(writer, sheet_name="Metadados", index=False)
        idata.to_excel(writer, sheet_name="Dados", index=False)

    treat_config = TratamentoConfig(
        input=input_path, imeta_sheet="Metadados", idata_sheet="Dados", sigla="TESTE", output_dir=tmp_path / "OUTPUT"
    )
    result = run_tratamento(treat_config)

    outfile = slides_resultT(
        template=_TEMPLATE,
        setor_estrategico="Teste Setor",
        sigla="TESTE",
        caminho_shp_mun=_SHP_MUN,
        caminho_shp_uf=_SHP_UF,
        resumo=result.resumo,
        meta_adapta=result.imeta,
        reference=result.ref,
        databruto=result.dados_b,
        datawinz=result.data_win,
        databxcx=result.data_bxc,
        datanorm=result.data_normal,
        output_dir=tmp_path / "OUTPUT",
    )

    assert outfile.exists()
    prs = Presentation(str(outfile))
    slides = list(prs.slides)
    # 2 header slides (subtitle + diagram) + 4 slides per indicator x 2 indicators
    assert len(slides) == 2 + 4 * 2
    assert slides[2].shapes.title.text == "Indicador Um - IND1"
    assert slides[6].shapes.title.text == "Indicador Dois - IND2"
